import pandas as pd
import openpyxl
import os
from pptx import Presentation
import comtypes.client
import pptxtopdf
from pathlib import Path
import shutil


def obter_dados_usuario():
    texto_do_curso = str(input('Digite qual foi o evento e carga horária. Dica de Texto: Concluiu o curso de... com uma carga horária de.... horas: '))
    data_emissao = input('Digite o local e a data seguindo este formato: Cidade, 01 de Janeiro de 2000: ')
    pasta = input('Digite o nome da pasta para salvar: ')
    return texto_do_curso, data_emissao, pasta

def criar_pasta(pasta):
    os.makedirs(pasta, exist_ok=True)

def carregar_dados_excel(caminho_arquivo):
    return pd.read_excel(caminho_arquivo)

def gerar_certificados(df, texto_do_curso, data_emissao, pasta_pptx, caminho_template):
    for _, row in df.iterrows():
        nome_completo = row['Nome Completo']
        template = Presentation(caminho_template)
        slide = template.slides[0]

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text == 'NOME COMPLETO':
                            run.text = nome_completo
                        if run.text == 'LOCAL E DATA.':
                            run.text = data_emissao
                        if run.text == 'CURSO E DURAÇÃO':
                            run.text = texto_do_curso

        nome_arquivo = f'{nome_completo}.pptx'
        caminho_completo = os.path.join(pasta_pptx, nome_arquivo)
        template.save(caminho_completo)

def converter_para_pdf(pasta_pptx, pasta_pdf):
    os.makedirs(pasta_pdf, exist_ok=True)
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    Arquivos_pptx = os.listdir(pasta_pptx)
    lista_pptx = [f for f in Arquivos_pptx if f.endswith(".pptx")]

    for pptx in lista_pptx:
        fullpath = os.path.join(pasta_pptx, pptx)
        pdf_filename = os.path.splitext(pptx)[0] + ".pdf"
        pdf_path = os.path.join(pasta_pdf, pdf_filename)

        try:
            pptxtopdf.convert(fullpath, pasta_pdf)
        except Exception as e:
            print(f"Erro ao converter {lista_pptx} usando pptxtopdf: {e}")
        print(f'{pdf_filename} criado com sucesso!')

def limpar_pasta(pasta):
    pasta = Path(pasta)
    for item in pasta.iterdir():
        if item.is_file():
            item.unlink()
        elif item.is_dir():
            shutil.rmtree(item)
    pasta.rmdir()

def main():
    texto_do_curso, data_emissao, pasta = obter_dados_usuario()
    pasta_pptx = rf"C:CAMINHO DO PROJETO\{pasta}"
    criar_pasta(pasta_pptx)

    df = carregar_dados_excel('ListaContato.xlsx')
    gerar_certificados(df, texto_do_curso, data_emissao, pasta_pptx,
                       r'C:CAMINHO DO CERTIFICADO PPTX')

    pasta_pdf = rf"C:CAMINHO DO PROJETO\{pasta}_PDF"
    converter_para_pdf(pasta_pptx, pasta_pdf)

    limpar_pasta(pasta_pptx)
    print("Conversão e organização dos certificados concluídas com sucesso!")

if __name__ == "__main__":
    main()
